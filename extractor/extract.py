"""
SharePoint Knowledge Base Extractor
------------------------------------
Authenticates via interactive browser SSO (opens your AB InBev login page),
reads all documents from a SharePoint document library, extracts text content,
and writes a structured documents.json for the React knowledge base app.

Usage:
    python extract.py
    # or double-click: run_local.ps1 / run_local.bat

A browser window opens for you to sign in with your AB InBev account.
The token is cached so subsequent runs are silent (no browser needed).
"""

import io
import json
import os
import re
import sys

# Force UTF-8 output on Windows (avoids encoding errors)
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')

import time
from datetime import datetime
from pathlib import Path
from typing import Optional

import msal
import requests

# ---------------------------------------------------------------------------
# Configuration — overridable via environment variables (used by GitHub Actions)
# ---------------------------------------------------------------------------

SHAREPOINT_HOST = os.environ.get("SHAREPOINT_HOST", "anheuserbuschinbev.sharepoint.com")
SITE_PATH       = os.environ.get("SITE_PATH",       "/sites/Sustinability")
DRIVE_NAME      = os.environ.get("DRIVE_NAME",      "Shared Documents")

# Azure AD public client — interactive browser SSO (no app registration needed)
# Uses Microsoft's well-known SharePoint Online Management Shell client ID
_PUBLIC_CLIENT_ID = "9bc3ab49-b65d-410a-85ad-de819febfddc"
AUTHORITY         = "https://login.microsoftonline.com/organizations"
SCOPES_DELEGATED  = ["https://graph.microsoft.com/Sites.Read.All",
                      "https://graph.microsoft.com/Files.Read.All"]
SCOPES_APP        = ["https://graph.microsoft.com/.default"]

OUTPUT_DIR  = Path(__file__).parent.parent / "knowledge-base-app" / "public" / "data"
OUTPUT_FILE = OUTPUT_DIR / "documents.json"
META_FILE   = OUTPUT_DIR / "kb-meta.json"

EXTRACTABLE = {".docx", ".pdf", ".pptx", ".xlsx", ".txt", ".md", ".csv"}
MEDIA_TYPES = {".mp4", ".mkv", ".avi", ".mov", ".mp3", ".wav",
               ".m4v", ".wmv", ".webm", ".m4a"}

# ---------------------------------------------------------------------------
# Auth helpers
# ---------------------------------------------------------------------------

def _is_ci() -> bool:
    """True when running in GitHub Actions with app credentials."""
    return all(os.environ.get(k) for k in ("TENANT_ID", "CLIENT_ID", "CLIENT_SECRET"))


def _get_token_client_credentials() -> str:
    """App-only auth via client credentials (used in GitHub Actions CI)."""
    tenant_id     = os.environ["TENANT_ID"]
    client_id     = os.environ["CLIENT_ID"]
    client_secret = os.environ["CLIENT_SECRET"]

    app = msal.ConfidentialClientApplication(
        client_id,
        authority=f"https://login.microsoftonline.com/{tenant_id}",
        client_credential=client_secret,
    )
    result = app.acquire_token_for_client(scopes=SCOPES_APP)
    if "access_token" not in result:
        raise RuntimeError(
            f"CI auth failed: {result.get('error')} — {result.get('error_description')}"
        )
    print("[AUTH] App-only token acquired via client credentials.", flush=True)
    return result["access_token"]


def _get_token() -> str:
    """Acquire a Graph API token — auto-selects CI vs interactive."""
    if _is_ci():
        return _get_token_client_credentials()
    return _get_token_interactive()


def _get_token_interactive() -> str:
    """Interactive browser SSO — opens your AB InBev login page in the browser."""
    cache = msal.SerializableTokenCache()
    cache_path = Path(__file__).parent / ".token_cache.json"
    if cache_path.exists():
        cache.deserialize(cache_path.read_text())

    app = msal.PublicClientApplication(
        _PUBLIC_CLIENT_ID, authority=AUTHORITY, token_cache=cache
    )

    # Try silent (cached token) first — skips browser on repeat runs
    accounts = app.get_accounts()
    if accounts:
        result = app.acquire_token_silent(SCOPES_DELEGATED, account=accounts[0])
        if result and "access_token" in result:
            cache_path.write_text(cache.serialize())
            print("[AUTH] Using cached token (no sign-in needed).", flush=True)
            return result["access_token"]

    # Interactive browser flow — opens a real browser window for SSO sign-in
    print("\nOpening browser for AB InBev SSO sign-in...", flush=True)
    print("If no window appears, check your taskbar.\n", flush=True)

    result = app.acquire_token_interactive(scopes=SCOPES_DELEGATED)

    if "access_token" not in result:
        # Fall back to device code if interactive not supported in this environment
        print("[WARN] Interactive auth failed, trying device code flow...", flush=True)
        flow = app.initiate_device_flow(scopes=SCOPES_DELEGATED)
        if "user_code" not in flow:
            raise RuntimeError(f"Auth failed: {flow.get('error_description')}")
        print("\n" + "=" * 60, flush=True)
        print("ACTION REQUIRED -- Sign in to Microsoft", flush=True)
        print("=" * 60, flush=True)
        print(flow["message"], flush=True)
        print("=" * 60 + "\n", flush=True)
        result = app.acquire_token_by_device_flow(flow)
        if "access_token" not in result:
            raise RuntimeError(f"Auth failed: {result.get('error_description')}")

    cache_path.write_text(cache.serialize())
    print("[OK]  Authenticated successfully!\n", flush=True)
    return result["access_token"]


# ---------------------------------------------------------------------------
# Graph API helpers
# ---------------------------------------------------------------------------

class GraphClient:
    BASE = "https://graph.microsoft.com/v1.0"

    def __init__(self, token: str):
        self.session = requests.Session()
        self.session.headers.update({"Authorization": f"Bearer {token}"})

    def get(self, path: str, **kwargs) -> dict:
        resp = self.session.get(f"{self.BASE}{path}", **kwargs)
        resp.raise_for_status()
        return resp.json()

    def get_bytes(self, url: str) -> bytes:
        resp = self.session.get(url)
        resp.raise_for_status()
        return resp.content

    def paginate(self, path: str) -> list:
        """Collect all pages from a paged Graph API response."""
        results = []
        url = f"{self.BASE}{path}"
        while url:
            resp = self.session.get(url)
            resp.raise_for_status()
            data = resp.json()
            results.extend(data.get("value", []))
            url = data.get("@odata.nextLink")
        return results


# ---------------------------------------------------------------------------
# SharePoint navigation
# ---------------------------------------------------------------------------

def get_site_id(client: GraphClient) -> str:
    data = client.get(f"/sites/{SHAREPOINT_HOST}:{SITE_PATH}")
    return data["id"]


def get_drive_id(client: GraphClient, site_id: str) -> str:
    drives = client.paginate(f"/sites/{site_id}/drives")
    for d in drives:
        if d["name"].lower() in (DRIVE_NAME.lower(), "documents", "shared documents"):
            return d["id"]
    # Fall back to default drive
    data = client.get(f"/sites/{site_id}/drive")
    return data["id"]


def list_items_recursive(client: GraphClient, drive_id: str,
                         folder_id: str = "root",
                         folder_path: str = "") -> list:
    """Recursively list all drive items under a folder."""
    items = []
    children = client.paginate(f"/drives/{drive_id}/items/{folder_id}/children")
    for item in children:
        item["_folderPath"] = folder_path
        if "folder" in item:
            sub_path = f"{folder_path}/{item['name']}".lstrip("/")
            items += list_items_recursive(client, drive_id, item["id"], sub_path)
        else:
            items.append(item)
    return items


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------

def extract_docx(content: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_pdf(content: bytes) -> str:
    import pdfplumber
    text_parts = []
    with pdfplumber.open(io.BytesIO(content)) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                text_parts.append(t)
    return "\n".join(text_parts)


def extract_pptx(content: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t:
                        slide_texts.append(t)
        if slide_texts:
            parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def extract_xlsx(content: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            row_str = "\t".join(str(c) if c is not None else "" for c in row)
            if row_str.strip():
                rows.append(row_str)
        if rows:
            parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows[:200]))
    return "\n\n".join(parts)


def extract_text(content: bytes, extension: str) -> str:
    try:
        if extension == ".docx":
            return extract_docx(content)
        elif extension == ".pdf":
            return extract_pdf(content)
        elif extension == ".pptx":
            return extract_pptx(content)
        elif extension == ".xlsx":
            return extract_xlsx(content)
        elif extension in (".txt", ".md", ".csv"):
            return content.decode("utf-8", errors="replace")
    except Exception as exc:
        return f"[Could not extract text: {exc}]"
    return ""


# ---------------------------------------------------------------------------
# Category / tag inference
# ---------------------------------------------------------------------------

CATEGORY_RULES = [
    (r"\.(mp4|mkv|avi|mov|webm|m4v)$", "Recordings"),
    (r"\.pptx?$", "Presentations"),
    (r"brd|business[\s_-]?requirement", "BRD"),
    (r"sow|statement[\s_-]?of[\s_-]?work", "SOW"),
    (r"meeting|minutes|mom", "Meeting Notes"),
    (r"report|dashboard", "Reports"),
    (r"template", "Templates"),
    (r"training|onboarding|guide|manual|tutorial", "Training"),
    (r"architecture|design|diagram|flow", "Architecture"),
    (r"\.xlsx?$", "Spreadsheets"),
    (r"\.pdf$", "Documents"),
    (r"\.docx?$", "Documents"),
]

# Maps file/folder names + content keywords → product IDs in the KB
PRODUCT_RULES = [
    ("circular-plan",     r"circular|retpack|ret[\s_-]?pack|capex.*pack|pack.*capex"),
    ("demand-cockpit",    r"demand[\s_-]?cockpit|demand[\s_-]?planning[\s_-]?cockpit"),
    ("material-planning", r"material[\s_-]?planning|e2e[\s_-]?material|\bmrp\b|on[\s_-]?time[\s_-]?suggest"),
    ("o9-adoption",       r"o9[\s_-]?adoption|touchless[\s_-]?plan|user[\s_-]?adoption|churn.*o9"),
    ("core-design",       r"core[\s_-]?design|core[\s_-]?process|as[\s_-]?is|process[\s_-]?map|catalogue"),
    ("o2d",               r"o2d|order[\s_-]?to[\s_-]?deliver|order[\s_-]?management|\bsto\b|\bstr\b|transport[\s_-]?schedul"),
]

def infer_product(name: str, folder_path: str, snippet: str = "") -> Optional[str]:
    combined = f"{folder_path}/{name} {snippet[:300]}".lower()
    for product_id, pattern in PRODUCT_RULES:
        if re.search(pattern, combined, re.IGNORECASE):
            return product_id
    return None

def infer_category(name: str, folder_path: str) -> str:
    combined = f"{folder_path}/{name}".lower()
    for pattern, category in CATEGORY_RULES:
        if re.search(pattern, combined, re.IGNORECASE):
            return category
    return "General"


def infer_tags(name: str, folder_path: str, text_snippet: str) -> list[str]:
    tags = set()
    combined = f"{folder_path}/{name} {text_snippet[:500]}".lower()
    tag_map = {
        "ai":             ["artificial intelligence", "machine learning", "llm", "copilot"],
        "sustainability": ["sustain", "esg", "environment", "retpack", "circular"],
        "gcc":            ["gcc", "global capability", "command centre"],
        "data":           ["data", "analytics", "dashboard", "kpi", "metric"],
        "process":        ["process", "workflow", "sop", "procedure", "as-is"],
        "finance":        ["budget", "cost", "finance", "revenue", "capex"],
        "planning":       ["planning", "forecast", "demand", "supply", "mrp"],
        "o9":             ["o9", "one planning", "oneplanning"],
        "project":        ["project", "milestone", "timeline", "roadmap"],
    }
    for tag, keywords in tag_map.items():
        if any(kw in combined for kw in keywords):
            tags.add(tag)
    return sorted(tags)


# ---------------------------------------------------------------------------
# Main extraction pipeline
# ---------------------------------------------------------------------------

def file_icon(ext: str) -> str:
    icons = {
        ".pdf": "📄", ".docx": "📝", ".doc": "📝",
        ".pptx": "📊", ".ppt": "📊", ".xlsx": "📈", ".xls": "📈",
        ".mp4": "🎥", ".mkv": "🎥", ".avi": "🎥", ".mov": "🎥",
        ".mp3": "🎵", ".wav": "🎵", ".m4a": "🎵",
        ".txt": "📃", ".md": "📃", ".csv": "📋",
        ".zip": "📦", ".rar": "📦",
    }
    return icons.get(ext.lower(), "📁")


def run():
    print("[AUTH]  Acquiring Microsoft Graph token...", flush=True)
    token = _get_token()
    client = GraphClient(token)

    print("[SITE]  Locating SharePoint site...", flush=True)
    site_id = get_site_id(client)
    print(f"    Site ID: {site_id}", flush=True)

    print("[DRIVE] Locating document library...", flush=True)
    drive_id = get_drive_id(client, site_id)
    print(f"    Drive ID: {drive_id}", flush=True)

    print("[LIST]  Listing all files (this may take a moment)...", flush=True)
    raw_items = list_items_recursive(client, drive_id)
    print(f"    Found {len(raw_items)} files", flush=True)

    documents = []
    total = len(raw_items)

    for idx, item in enumerate(raw_items, 1):
        name = item.get("name", "")
        ext = Path(name).suffix.lower()
        size = item.get("size", 0)
        folder_path = item.get("_folderPath", "")
        modified = item.get("lastModifiedDateTime", "")
        created = item.get("createdDateTime", "")
        web_url = item.get("webUrl", "")
        author = (item.get("createdBy", {}).get("user", {}).get("displayName")
                  or item.get("lastModifiedBy", {}).get("user", {}).get("displayName")
                  or "Unknown")

        print(f"  [{idx}/{total}] {name}", end="", flush=True)

        text_content = ""
        is_media = ext in MEDIA_TYPES

        if not is_media and ext in EXTRACTABLE and size < 50 * 1024 * 1024:  # skip >50MB
            try:
                download_url = item.get("@microsoft.graph.downloadUrl") or \
                               item.get("microsoft.graph.downloadUrl")
                if not download_url:
                    # Fetch fresh item with download URL
                    fresh = client.get(f"/drives/{drive_id}/items/{item['id']}")
                    download_url = fresh.get("@microsoft.graph.downloadUrl")

                if download_url:
                    content = client.get_bytes(download_url)
                    text_content = extract_text(content, ext)
                    print(f" [OK] ({len(text_content)} chars)", flush=True)
                else:
                    print(" [WARN] no download URL", flush=True)
            except Exception as exc:
                print(f" [ERR] {exc}", flush=True)
        elif is_media:
            print(" [MEDIA] skipped extraction", flush=True)
        else:
            print(" [SKIP]", flush=True)

        snippet = " ".join(text_content.split()[:80]) if text_content else ""

        documents.append({
            "id": item["id"],
            "name": name,
            "extension": ext,
            "icon": file_icon(ext),
            "size": size,
            "sizeLabel": _format_size(size),
            "folderPath": folder_path,
            "category": infer_category(name, folder_path),
            "productId": infer_product(name, folder_path, snippet),
            "tags": infer_tags(name, folder_path, snippet),
            "author": author,
            "modified": modified,
            "created": created,
            "webUrl": web_url,
            "isMedia": is_media,
            "snippet": snippet,
            "content": text_content,
        })

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    with open(OUTPUT_FILE, "w", encoding="utf-8") as f:
        json.dump({
            "generated": datetime.utcnow().isoformat() + "Z",
            "totalDocuments": len(documents),
            "documents": documents,
        }, f, ensure_ascii=False, indent=2)

    print(f"\n[DONE]  {len(documents)} documents indexed --> {OUTPUT_FILE}", flush=True)
    print("    Now run:  cd knowledge-base-app && npm run dev", flush=True)


def _format_size(size: int) -> str:
    for unit in ("B", "KB", "MB", "GB"):
        if size < 1024:
            return f"{size:.1f} {unit}"
        size /= 1024
    return f"{size:.1f} TB"


if __name__ == "__main__":
    run()
